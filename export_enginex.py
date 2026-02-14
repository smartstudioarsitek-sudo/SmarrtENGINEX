import pandas as pd
import io
import re
import folium
from streamlit_folium import st_folium
import json
import xml.etree.ElementTree as ET

# Library untuk Dokumen Office
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
import xlsxwriter
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt

class EnginexExporter:
    """
    Modul Khusus ENGINEX Ultimate untuk:
    1. BIM/CAD Data Exchange (CSV/TXT)
    2. Geospatial Visualization (Folium)
    3. Formal Reporting (PUPR Standard - Excel, Word, PPT)
    """

    # ==========================================
    # 1. BIM & CIVIL 3D INTEGRATION
    # ==========================================
    @staticmethod
    def export_to_civil3d_csv(data_list):
        """
        Mengubah data list/dataframe menjadi format CSV untuk Import Points Civil 3D.
        Format Standar: P,N,E,Z,D (Point#, Northing, Easting, Elevation, Description)
        """
        try:
            # Asumsi data masuk berupa DataFrame atau List of Dicts
            if isinstance(data_list, list):
                df = pd.DataFrame(data_list)
            else:
                df = data_list

            # Validasi kolom minimal (harus ada koordinat)
            required_cols = ['x', 'y', 'z'] # Easting, Northing, Elevation
            if not all(col in df.columns.str.lower() for col in required_cols):
                # Jika kolom tidak sesuai, return None atau CSV biasa
                return df.to_csv(index=False).encode('utf-8')

            # Mapping ke standar PNEZD
            # Pastikan urutan kolom: PointNumber, Northing(Y), Easting(X), Elevation(Z), Desc
            output_df = pd.DataFrame()
            output_df['P'] = range(1, len(df) + 1)
            output_df['N'] = df.get('y', df.get('Y', 0))
            output_df['E'] = df.get('x', df.get('X', 0))
            output_df['Z'] = df.get('z', df.get('Z', 0))
            output_df['D'] = df.get('desc', df.get('keterangan', 'Point'))

            return output_df.to_csv(index=False, header=False).encode('utf-8')
        except Exception as e:
            return f"Error BIM Export: {str(e)}"

    # ==========================================
    # 2. GEOSPATIAL VISUALIZATION (FOLIUM)
    # ==========================================
    @staticmethod
    def render_geospatial_map(file_content, file_type):
        """
        Merender peta Folium dari GeoJSON atau KML (Parsing XML sederhana).
        Return: Objek Peta Folium
        """
        m = folium.Map(location=[-2.5, 118.0], zoom_start=5, tiles="OpenStreetMap") # Default Indonesia
        
        try:
            if file_type == 'geojson':
                data = json.loads(file_content)
                folium.GeoJson(data, name="GeoJSON Layer").add_to(m)
                
                # Auto-zoom ke fitur pertama
                if data['features']:
                    coords = data['features'][0]['geometry']['coordinates']
                    # Simplifikasi pengambilan centroid untuk zoom (bisa dikembangkan)
                    if len(coords) > 0:
                        # Logic sederhana untuk point/polygon
                        pass 

            elif file_type in ['kml', 'kmz']:
                # KML Parsing Sederhana (Extract Coordinates via XML)
                # Note: Full KML support butuh library berat (GDAL/fiona) yg sulit di Streamlit Cloud
                # Ini pendekatan "Lightweight"
                root = ET.fromstring(file_content)
                namespace = {'kml': 'http://www.opengis.net/kml/2.2'}
                
                points = []
                for placemark in root.findall('.//kml:Placemark', namespace):
                    name = placemark.find('kml:name', namespace).text if placemark.find('kml:name', namespace) is not None else "No Name"
                    coords_str = placemark.find('.//kml:coordinates', namespace)
                    
                    if coords_str is not None:
                        coords = coords_str.text.strip().split()
                        for c in coords:
                            lon, lat, *_ = c.split(',')
                            points.append([float(lat), float(lon)])
                            folium.Marker([float(lat), float(lon)], popup=name).add_to(m)
                
                if points:
                    m.fit_bounds(points)

            folium.LayerControl().add_to(m)
            return m
        except Exception as e:
            return None

    # ==========================================
    # 3. PUPR EXCEL EXPORT (RAB STYLE)
    # ==========================================
    @staticmethod
    def create_pupr_excel(dataframe, sheet_name="Data Teknis"):
        """
        Membuat Excel dengan format Tabel Teknik (Border, Header Bold, Auto-width).
        """
        output = io.BytesIO()
        workbook = xlsxwriter.Workbook(output, {'in_memory': True})
        worksheet = workbook.add_worksheet(sheet_name)

        # Format Styles
        header_format = workbook.add_format({
            'bold': True,
            'text_wrap': True,
            'valign': 'vcenter',
            'align': 'center',
            'fg_color': '#D7E4BC', # Hijau muda standar excel
            'border': 1
        })
        
        cell_format = workbook.add_format({
            'border': 1,
            'valign': 'top'
        })
        
        number_format = workbook.add_format({
            'border': 1,
            'valign': 'top',
            'num_format': '#,##0.00'
        })

        # Write Header
        for col_num, value in enumerate(dataframe.columns.values):
            worksheet.write(0, col_num, value, header_format)
            worksheet.set_column(col_num, col_num, 20) # Default width

        # Write Data
        for row_num, row_data in enumerate(dataframe.values):
            for col_num, cell_value in enumerate(row_data):
                # Cek tipe data untuk formatting angka
                if isinstance(cell_value, (int, float)):
                    worksheet.write(row_num + 1, col_num, cell_value, number_format)
                else:
                    worksheet.write(row_num + 1, col_num, str(cell_value), cell_format)

        workbook.close()
        output.seek(0)
        return output

    # ==========================================
    # 4. PUPR WORD EXPORT (LAPORAN RESMI)
    # ==========================================
    @staticmethod
    def create_pupr_word(text_content, project_name="Proyek ENGINEX"):
        """
        Membuat Dokumen Word dengan style Laporan Teknik.
        Menerjemahkan Markdown sederhana ke Word formatting.
        """
        doc = Document()
        
        # --- SETTING KOP SURAT (Placeholder) ---
        header = doc.sections[0].header
        htable = header.add_table(1, 2, width=Inches(6))
        htable.autofit = False
        htable.columns[0].width = Inches(1)
        htable.columns[1].width = Inches(5)
        
        # Kolom 1: Logo/Nama (Simulasi)
        cell_logo = htable.cell(0, 0)
        cell_logo.text = "ENGINEX"
        
        # Kolom 2: Judul Proyek
        cell_text = htable.cell(0, 1)
        p = cell_text.paragraphs[0]
        p.text = f"LAPORAN TEKNIS: {project_name}\nGenerated by AI Consultant"
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT

        # --- CONTENT PARSING ---
        doc.add_heading(f'Laporan Analisis: {project_name}', 0)
        
        lines = text_content.split('\n')
        for line in lines:
            clean_line = line.strip()
            
            # Heading Levels
            if clean_line.startswith('### '):
                doc.add_heading(clean_line.replace('### ', ''), level=2)
            elif clean_line.startswith('## '):
                doc.add_heading(clean_line.replace('## ', ''), level=1)
            
            # Bullet Points
            elif clean_line.startswith('- ') or clean_line.startswith('* '):
                p = doc.add_paragraph(clean_line[2:], style='List Bullet')
            
            # Numbered List (Sederhana)
            elif re.match(r'^\d+\.', clean_line):
                p = doc.add_paragraph(clean_line, style='List Number')
                
            # Regular Text
            elif clean_line:
                p = doc.add_paragraph(clean_line)
                p.paragraph_format.space_after = Pt(6)

        bio = io.BytesIO()
        doc.save(bio)
        bio.seek(0)
        return bio

    # ==========================================
    # 5. PUPR POWERPOINT EXPORT
    # ==========================================
    @staticmethod
    def create_pupr_pptx(text_content, project_name):
        """
        Membuat Slide Presentasi dari ringkasan konten.
        """
        prs = Presentation()
        
        # Slide 1: Title
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = project_name
        subtitle.text = "Laporan Analisis Teknis & Rekomendasi\nDisusun oleh ENGINEX Ultimate"

        # Parsing konten untuk slide content
        # Logika: Setiap Heading (##) menjadi Judul Slide Baru
        current_slide = None
        body_text_frame = None
        
        lines = text_content.split('\n')
        for line in lines:
            clean_line = line.strip()
            
            if clean_line.startswith('## ') or clean_line.startswith('### '):
                # Buat Slide Baru
                bullet_slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = current_slide.shapes
                title_shape = shapes.title
                title_shape.text = clean_line.replace('#', '').strip()
                body_text_frame = shapes.placeholders[1].text_frame
                
            elif clean_line and body_text_frame:
                # Isi konten slide
                p = body_text_frame.add_paragraph()
                p.text = clean_line.replace('*', '').replace('-', '').strip()
                p.level = 0
                if clean_line.startswith('  '): # Indentasi sederhana
                    p.level = 1

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output
