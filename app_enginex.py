import streamlit as st
import google.generativeai as genai
from google.generativeai.types import HarmCategory, HarmBlockThreshold
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import json
from PIL import Image
import PyPDF2
import io
import docx
from docx.shared import Inches, Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
import zipfile
from pptx import Presentation
import xlsxwriter
import re
import time

# ==========================================
# 0. CEK MODUL PENDUKUNG (BACKEND & PERSONA)
# ==========================================
try:
    from backend_enginex import EnginexBackend
    from persona import gems_persona, get_persona_list, get_system_instruction
    from streamlit_folium import st_folium
except ImportError as e:
    st.error(f"❌ CRITICAL ERROR: Modul pendukung tidak ditemukan! \nDetail Error: {e}")
    st.warning("Pastikan file 'backend_enginex.py' dan 'persona.py' sudah ada di folder yang sama.")
    st.stop()

# ==========================================
# 1. MESIN PEMBUAT LAPORAN (REPORT ENGINE)
# ==========================================
class UniversalReportGenerator:
    """
    Kelas ini menangani pembuatan semua dokumen (Word, Excel, PPT, Zip).
    Ditanam langsung di sini agar integrasi data visual berjalan mulus.
    """
    
    @staticmethod
    def create_word_report_pupr(project_name, expert_name, text_content, visual_assets):
        """
        Membuat Laporan Word Standar PU-PR:
        - Kop Surat
        - Format Teks Rapi
        - Lampiran Gambar & Grafik Otomatis
        """
        doc = docx.Document()
        
        # --- A. SETUP HALAMAN (A4 & Margin 2.54cm) ---
        section = doc.sections[0]
        section.page_height = Cm(29.7)
        section.page_width = Cm(21.0)
        section.left_margin = Cm(2.54)
        section.right_margin = Cm(2.54)
        section.top_margin = Cm(2.54)
        section.bottom_margin = Cm(2.54)
        
        # --- B. HEADER / KOP SURAT ---
        header = section.header
        htable = header.add_table(1, 2, width=Inches(6))
        htable.autofit = False
        htable.columns[0].width = Inches(2.5)
        htable.columns[1].width = Inches(3.5)
        
        # Sel Kiri: Identitas Konsultan
        cell_logo = htable.cell(0, 0)
        cell_logo.text = "ENGINEX CONSULTANT\nIntegrated AI System"
        
        # Sel Kanan: Info Proyek
        cell_info = htable.cell(0, 1)
        p_head = cell_info.paragraphs[0]
        p_head.text = f"PROYEK: {project_name}\nAHLI: {expert_name}\nTANGGAL: {time.strftime('%d-%m-%Y')}"
        p_head.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        
        # --- C. JUDUL UTAMA ---
        doc.add_paragraph("\n") # Spasi
        title = doc.add_heading("LAPORAN ANALISIS TEKNIS", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("-" * 70).alignment = WD_ALIGN_PARAGRAPH.CENTER

        # --- D. PARSING KONTEN TEKS (MARKDOWN TO WORD) ---
        lines = text_content.split('\n')
        for line in lines:
            line = line.strip()
            if not line: continue
            
            # Deteksi Heading Markdown
            if line.startswith('### '):
                doc.add_heading(line.replace('### ', ''), level=3)
            elif line.startswith('## '):
                doc.add_heading(line.replace('## ', ''), level=2)
            elif line.startswith('# '):
                doc.add_heading(line.replace('# ', ''), level=1)
            # Deteksi List Bullet
            elif line.startswith('- ') or line.startswith('* '):
                doc.add_paragraph(line[2:], style='List Bullet')
            # Deteksi List Angka
            elif re.match(r'^\d+\.', line):
                doc.add_paragraph(line, style='List Number')
            # Paragraf Biasa
            else:
                p = doc.add_paragraph(line)
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

        # --- E. LAMPIRAN VISUAL (GRAFIK & GAMBAR) ---
        if visual_assets:
            doc.add_page_break()
            doc.add_heading("LAMPIRAN: DOKUMENTASI VISUAL & GRAFIK", level=1)
            doc.add_paragraph("Berikut adalah hasil render grafik teknis dan gambar pendukung:")
            
            for i, asset in enumerate(visual_assets):
                try:
                    # Tambahkan Caption
                    caption_text = asset.get('caption', f'Gambar-{i+1}')
                    doc.add_heading(f"Gambar {i+1}: {caption_text}", level=3)
                    
                    # Proses Gambar (Dari Memory BytesIO)
                    img_stream = asset['data']
                    
                    # Kita harus copy stream agar tidak corrupt jika dipakai ulang
                    temp_stream = io.BytesIO()
                    
                    if asset['type'] == 'image':
                        # Jika format PIL Image
                        asset['data'].save(temp_stream, format='PNG')
                    else:
                        # Jika format Plot BytesIO
                        img_stream.seek(0)
                        temp_stream.write(img_stream.read())
                    
                    temp_stream.seek(0)
                    doc.add_picture(temp_stream, width=Inches(6.0))
                    doc.add_paragraph("") # Spasi kosong bawah gambar
                    
                except Exception as e:
                    doc.add_paragraph(f"[Error memuat gambar {i+1}: {str(e)}]")

        # --- F. FOOTER ---
        footer_section = doc.sections[0]
        footer = footer_section.footer
        p_foot = footer.paragraphs[0]
        p_foot.text = "Dokumen ini digenerate secara otomatis oleh ENGINEX AI."
        p_foot.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Simpan ke Buffer
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer

    @staticmethod
    def create_excel_table(text_content):
        """Mendeteksi tabel markdown dan mengubahnya jadi Excel"""
        try:
            lines = text_content.split('\n')
            data_rows = []
            headers = []
            
            for line in lines:
                # Cek format tabel markdown | A | B |
                if "|" in line and "---" not in line:
                    row = [c.strip() for c in line.split('|')[1:-1]]
                    if not headers:
                        headers = row
                    else:
                        if len(row) == len(headers): # Pastikan jumlah kolom sama
                            data_rows.append(row)
            
            if not headers or not data_rows:
                return None
            
            # Buat DataFrame
            df = pd.DataFrame(data_rows, columns=headers)
            
            # Write to Excel Buffer
            output = io.BytesIO()
            with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                df.to_excel(writer, index=False, sheet_name='Data Analisis')
                # Auto-adjust column width
                worksheet = writer.sheets['Data Analisis']
                for idx, col in enumerate(df.columns):
                    worksheet.set_column(idx, idx, 20)
            
            output.seek(0)
            return output
        except:
            return None

    @staticmethod
    def create_zip_images(visual_assets):
        """Membungkus semua gambar jadi satu file ZIP"""
        if not visual_assets: return None
        
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w") as zf:
            for i, asset in enumerate(visual_assets):
                file_name = f"Gambar_{i+1}.png"
                
                img_data = io.BytesIO()
                if asset['type'] == 'image':
                    asset['data'].save(img_data, format='PNG')
                else: # plot
                    asset['data'].seek(0)
                    img_data.write(asset['data'].read())
                
                img_data.seek(0)
                zf.writestr(file_name, img_data.read())
                
        zip_buffer.seek(0)
        return zip_buffer

    @staticmethod
    def create_civil3d_csv(text_content):
        """Ekstrak koordinat P,N,E,Z,D untuk Civil 3D"""
        # Coba ambil tabel dulu
        excel_io = UniversalReportGenerator.create_excel_table(text_content)
        if not excel_io: return None
        
        try:
            df = pd.read_excel(excel_io)
            cols = [c.lower() for c in df.columns]
            
            # Syarat minimal: Ada X dan Y (atau Easting dan Northing)
            has_x = any('x' in c or 'east' in c for c in cols)
            has_y = any('y' in c or 'north' in c for c in cols)
            
            if has_x and has_y:
                # Mapping Kolom
                out = pd.DataFrame()
                out['P'] = range(1, len(df)+1)
                
                # Cari nama kolom asli
                col_x = next(c for c in df.columns if 'x' in c.lower() or 'east' in c.lower())
                col_y = next(c for c in df.columns if 'y' in c.lower() or 'north' in c.lower())
                col_z = next((c for c in df.columns if 'z' in c.lower() or 'elev' in c.lower()), None)
                col_d = next((c for c in df.columns if 'ket' in c.lower() or 'desc' in c.lower()), None)
                
                out['N'] = df[col_y]
                out['E'] = df[col_x]
                out['Z'] = df[col_z] if col_z else 0
                out['D'] = df[col_d] if col_d else "Point"
                
                return out.to_csv(index=False, header=False).encode('utf-8')
        except:
            pass
        return None

    @staticmethod
    def create_ppt_presentation(text_content, project_name):
        """Buat PPT sederhana"""
        prs = Presentation()
        # Slide Judul
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = project_name
        slide.placeholders[1].text = "Laporan Analisis ENGINEX AI"
        
        # Slide Isi
        lines = text_content.split('\n')
        curr_slide = None
        for line in lines:
            line = line.strip()
            if line.startswith('## '):
                curr_slide = prs.slides.add_slide(prs.slide_layouts[1])
                curr_slide.shapes.title.text = line.replace('## ', '')
            elif line.startswith('- ') and curr_slide:
                tf = curr_slide.shapes.placeholders[1].text_frame
                p = tf.add_paragraph()
                p.text = line.replace('- ', '')
                p.level = 0
        
        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return out

# ==========================================
# 2. KONFIGURASI HALAMAN & CSS
# ==========================================
st.set_page_config(
    page_title="ENGINEX Ultimate v15",
    page_icon="🏗️",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
<style>
    /* Header Utama */
    .main-header {
        font-size: 2.5rem;
        font-weight: 800;
        color: #1E293B;
        margin-bottom: 20px;
        border-bottom: 5px solid #2563EB;
        padding-bottom: 10px;
    }
    
    /* Tombol Download di Sidebar */
    .stDownloadButton button {
        width: 100%;
        background-color: #FFFFFF;
        color: #334155;
        border: 1px solid #94A3B8;
        font-weight: 600;
        text-align: left;
        padding: 10px 15px;
        margin-bottom: 8px;
        border-radius: 8px;
        transition: all 0.2s;
    }
    .stDownloadButton button:hover {
        background-color: #EFF6FF;
        border-color: #2563EB;
        color: #1D4ED8;
        transform: translateX(5px);
    }
    
    /* Header Kecil di Sidebar */
    .download-header {
        font-size: 0.9rem;
        font-weight: bold;
        color: #0F172A;
        margin-top: 15px;
        margin-bottom: 5px;
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }
    
    /* Kotak Pesan Auto-Pilot */
    .auto-pilot-box {
        background-color: #ECFDF5;
        border-left: 4px solid #10B981;
        padding: 10px;
        border-radius: 4px;
        color: #065F46;
        font-size: 0.9rem;
        margin-bottom: 10px;
    }
</style>
""", unsafe_allow_html=True)

# ==========================================
# 3. INITIALIZE SESSION STATE (MEMORY)
# ==========================================
if 'processed_files' not in st.session_state:
    st.session_state.processed_files = set()

if 'current_expert_active' not in st.session_state:
    st.session_state.current_expert_active = "👑 The GEMS Grandmaster"

if 'backend' not in st.session_state:
    st.session_state.backend = EnginexBackend()

# --- PENTING: MEMORY UNTUK TOMBOL DOWNLOAD ---
# Agar tombol tidak hilang saat refresh, kita simpan datanya di sini
if 'last_analysis_text' not in st.session_state:
    st.session_state.last_analysis_text = ""
    
if 'generated_visuals' not in st.session_state:
    st.session_state.generated_visuals = [] # List untuk menyimpan gambar/plot

db = st.session_state.backend

# ==========================================
# 4. FUNGSI UTILITAS (HELPER)
# ==========================================
def execute_generated_code(code_str):
    """
    Mengeksekusi kode Python (Matplotlib), menampilkan di layar,
    DAN MENYIMPANNYA ke memory untuk didownload.
    """
    try:
        # 1. Siapkan variabel
        local_vars = {"pd": pd, "np": np, "plt": plt, "st": st}
        plt.clf() # Bersihkan canvas lama
        
        # 2. Eksekusi Kode
        exec(code_str, {}, local_vars)
        
        # 3. Tangkap Gambar (Figure)
        fig = plt.gcf()
        
        # 4. Simpan ke Buffer (Memory)
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        buf.seek(0)
        
        # 5. Masukkan ke Session State (Agar bisa didownload)
        st.session_state.generated_visuals.append({
            'type': 'plot',
            'data': buf,
            'caption': "Grafik Analisis Teknis"
        })
        
        return True
    except Exception as e:
        st.error(f"⚠️ Gagal Render Grafik: {str(e)}")
        return False

def process_uploaded_file(uploaded_file):
    """Membaca berbagai jenis file upload"""
    if uploaded_file is None: return None, None
    file_type = uploaded_file.name.split('.')[-1].lower()
    
    try:
        if file_type in ['png', 'jpg', 'jpeg']:
            img = Image.open(uploaded_file)
            return "image", img
        elif file_type == 'pdf':
            pdf = PyPDF2.PdfReader(uploaded_file)
            text = "\n".join([p.extract_text() for p in pdf.pages])
            return "text", text
        elif file_type == 'docx':
            doc = docx.Document(uploaded_file)
            text = "\n".join([p.text for p in doc.paragraphs])
            return "text", text
        elif file_type in ['xlsx', 'xls']:
            df = pd.read_excel(uploaded_file)
            return "text", f"[DATA EXCEL PREVIEW]\n{df.head(50).to_string()}"
        elif file_type in ['geojson', 'kml']:
            return "geo", uploaded_file.getvalue().decode("utf-8")
        else:
            return "text", uploaded_file.getvalue().decode("utf-8")
    except Exception as e:
        return "error", str(e)

# ==========================================
# 5. SIDEBAR LAYOUT (DENGAN DOWNLOAD CENTER)
# ==========================================
with st.sidebar:
    st.title("🏗️ ENGINEX V15")
    st.caption("Ultimate Integrated System")
    st.markdown("---")
    
    # --- A. API KEY ---
    api_key_in = st.text_input("🔑 Google API Key:", type="password")
    api_key = api_key_in.strip() if api_key_in else st.secrets.get("GOOGLE_API_KEY", "")
    
    if not api_key:
        st.warning("⚠️ Masukkan API Key")
        st.stop()
    try:
        genai.configure(api_key=api_key)
    except:
        st.error("API Key Invalid")
        st.stop()

    # --- B. MODEL SELECTOR ---
    future_models = [
        "models/gemini-2.5-flash-lite",
        "models/gemini-2.5-flash-image",
        "models/gemini-1.5-pro",
        "models/gemini-1.5-flash"
    ]
    selected_model = st.selectbox("🧠 Model AI:", future_models, index=2)

    # --- C. MANAJEMEN PROYEK ---
    with st.expander("📁 Proyek & Ahli", expanded=False):
        projects = db.daftar_proyek()
        mode = st.radio("Mode:", ["Buka Lama", "Buat Baru"])
        
        if mode == "Buat Baru":
            nama_proyek = st.text_input("Nama Proyek:", "DED Bendungan 2026")
        else:
            nama_proyek = st.selectbox("Pilih Proyek:", projects) if projects else "Proyek Baru"
            
        st.markdown("---")
        auto_pilot = st.checkbox("🤖 Auto-Pilot", value=True)
        manual_expert = st.selectbox("Pilih Ahli:", get_persona_list(), disabled=auto_pilot)

    # --- D. UPLOAD FILE ---
    st.markdown("### 📎 Upload Data")
    uploaded_files = st.file_uploader("", accept_multiple_files=True)
    
    if st.button("🧹 Reset Sesi"):
        db.clear_chat(nama_proyek, st.session_state.current_expert_active)
        st.session_state.generated_visuals = [] # Bersihkan gambar
        st.session_state.last_analysis_text = "" # Bersihkan teks
        st.rerun()
    
    st.markdown("---")
    
    # ==========================================
    # 🌟 DOWNLOAD CENTER (FIXED & STABLE)
    # ==========================================
    st.header("📥 Download Center")
    
    # Kita cek variabel session state. Jika ada isi, tampilkan tombol.
    has_content = len(st.session_state.last_analysis_text) > 10
    
    if has_content:
        txt = st.session_state.last_analysis_text
        vis = st.session_state.generated_visuals
        
        # 1. LAPORAN WORD (TEKS + GAMBAR)
        st.markdown('<div class="download-header">📄 Laporan Resmi (Word)</div>', unsafe_allow_html=True)
        word_buffer = UniversalReportGenerator.create_word_report_pupr(
            nama_proyek, 
            st.session_state.current_expert_active, 
            txt, 
            vis
        )
        st.download_button(
            label="Download Laporan (.docx)",
            data=word_buffer,
            file_name=f"{nama_proyek}_Laporan_Lengkap.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        
        # 2. DATA TABEL (EXCEL)
        st.markdown('<div class="download-header">📊 Data Perhitungan (Excel)</div>', unsafe_allow_html=True)
        excel_buffer = UniversalReportGenerator.create_excel_table(txt)
        if excel_buffer:
            st.download_button(
                label="Download Tabel (.xlsx)",
                data=excel_buffer,
                file_name=f"{nama_proyek}_Data.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
        else:
            st.caption("ℹ️ Tidak ada tabel terdeteksi.")
            
        # 3. PRESENTASI (PPT)
        st.markdown('<div class="download-header">📢 Slide Presentasi (PPT)</div>', unsafe_allow_html=True)
        ppt_buffer = UniversalReportGenerator.create_ppt_presentation(txt, nama_proyek)
        st.download_button(
            label="Download Slide (.pptx)",
            data=ppt_buffer,
            file_name=f"{nama_proyek}_Presentasi.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
        # 4. ALBUM GAMBAR (ZIP)
        if vis:
            st.markdown('<div class="download-header">🖼️ Album Grafik (ZIP)</div>', unsafe_allow_html=True)
            zip_buffer = UniversalReportGenerator.create_zip_images(vis)
            st.download_button(
                label="Download Semua Gambar (.zip)",
                data=zip_buffer,
                file_name=f"{nama_proyek}_Images.zip",
                mime="application/zip"
            )
            
        # 5. CIVIL 3D (CSV)
        csv_buffer = UniversalReportGenerator.create_civil3d_csv(txt)
        if csv_buffer:
            st.markdown('<div class="download-header">🏗️ Civil 3D Points</div>', unsafe_allow_html=True)
            st.download_button(
                label="Download Points (.csv)",
                data=csv_buffer,
                file_name=f"{nama_proyek}_Points.csv",
                mime="text/csv"
            )
            
    else:
        st.info("Silakan mulai konsultasi agar menu download aktif.")

# ==========================================
# 6. LOGIKA CHAT UTAMA
# ==========================================
st.markdown(f'<div class="main-header">{nama_proyek}</div>', unsafe_allow_html=True)

# Tentukan Ahli Aktif
if auto_pilot:
    active_expert = st.session_state.current_expert_active
else:
    active_expert = manual_expert
    st.session_state.current_expert_active = manual_expert

st.caption(f"Status: **Connected** | Ahli Aktif: **{active_expert}**")

# --- LOAD HISTORY DARI DATABASE ---
chat_history = db.get_chat_history(nama_proyek, active_expert)

# Render Chat Lama
for chat in chat_history:
    with st.chat_message(chat['role']):
        st.markdown(chat['content'])
        
        # Render Grafik Arsip (Visual Only)
        if chat['role'] == 'assistant' and "plt." in chat['content']:
            with st.expander("📊 Lihat Grafik (Arsip)"):
                match = re.search(r"```python(.*?)```", chat['content'], re.DOTALL)
                if match:
                    try:
                        exec(match.group(1), {"pd":pd, "np":np, "plt":plt, "st":st})
                        st.pyplot(plt.gcf())
                        plt.clf()
                    except: pass

# --- INPUT USER ---
user_input = st.chat_input(f"Konsultasi dengan {active_expert}...")

if user_input:
    # 1. Routing Auto-Pilot
    target_expert = active_expert
    if auto_pilot:
        try:
            router = genai.GenerativeModel("models/gemini-1.5-flash")
            resp = router.generate_content(f"Pilih SATU ahli dari {list(gems_persona.keys())} untuk: '{user_input}'. Output Nama Saja.")
            suggestion = resp.text.strip()
            if suggestion in gems_persona:
                target_expert = suggestion
                st.session_state.current_expert_active = target_expert
        except: pass
    
    if target_expert != active_expert:
        st.markdown(f'<div class="auto-pilot-box">🤖 Auto-Pilot mengalihkan ke: <b>{target_expert}</b></div>', unsafe_allow_html=True)
        active_expert = target_expert

    # 2. Simpan & Tampilkan Pesan User
    db.simpan_chat(nama_proyek, active_expert, "user", user_input)
    with st.chat_message("user"):
        st.markdown(user_input)

    # 3. Persiapan Konten AI
    payload = [user_input]
    
    # (Opsional) Reset visual untuk chat baru agar tidak menumpuk dari sesi lama
    # st.session_state.generated_visuals = [] 
    
    if uploaded_files:
        for f in uploaded_files:
            if f.name not in st.session_state.processed_files:
                typ, cont = process_uploaded_file(f)
                if typ == "image":
                    with st.chat_message("user"): st.image(f, width=250)
                    # Simpan gambar upload ke daftar visual untuk laporan
                    buf = io.BytesIO()
                    cont.save(buf, format='PNG')
                    st.session_state.generated_visuals.append({'type':'image', 'data':buf, 'caption':f"Upload: {f.name}"})
                    payload.append(cont)
                elif typ == "text":
                    payload.append(f"\n[FILE: {f.name}]\n{cont}")
                st.session_state.processed_files.add(f.name)

    # 4. Generate AI Response
    with st.chat_message("assistant"):
        with st.spinner(f"{active_expert} sedang bekerja..."):
            try:
                # Instruksi Plotting
                sys_instr = get_system_instruction(active_expert)
                if "Code" not in active_expert and "Visionary" not in active_expert:
                    sys_instr += "\n[VISUALISASI]: Jika perlu grafik, buat kode Python (matplotlib). Akhiri dengan 'st.pyplot(plt.gcf())'."

                # Fallback Model Strategy
                try:
                    model = genai.GenerativeModel(selected_model, system_instruction=sys_instr)
                except:
                    model = genai.GenerativeModel("models/gemini-1.5-pro", system_instruction=sys_instr)

                # History Context
                hist_objs = []
                recent_chats = db.get_chat_history(nama_proyek, active_expert)[-5:]
                for h in recent_chats:
                    if h['content'] != user_input:
                        role = "user" if h['role']=="user" else "model"
                        hist_objs.append({"role": role, "parts": [h['content']]})

                # Start Chat
                chat = model.start_chat(history=hist_objs)
                
                # Handling Stream
                full_resp = ""
                box = st.empty()
                try:
                    stream = chat.send_message(payload, stream=True)
                    for chunk in stream:
                        if chunk.text:
                            full_resp += chunk.text
                            box.markdown(full_resp + "▌")
                except:
                    # Retry without stream (kadang model flash error di stream)
                    fallback_model = genai.GenerativeModel("models/gemini-1.5-pro", system_instruction=sys_instr)
                    full_resp = fallback_model.start_chat(history=hist_objs).send_message(payload).text
                
                # Final Display
                box.markdown(full_resp)
                db.simpan_chat(nama_proyek, active_expert, "assistant", full_resp)
                
                # 5. EKSEKUSI PLOT & SIMPAN KE MEMORY
                codes = re.findall(r"```python(.*?)```", full_resp, re.DOTALL)
                for c in codes:
                    if "plt." in c:
                        st.markdown("### 📉 Grafik Engineering")
                        with st.container():
                            # Fungsi ini otomatis menyimpan grafik ke st.session_state.generated_visuals
                            success = execute_generated_code(c)
                            if success:
                                st.pyplot(plt.gcf())
                                plt.clf()

                # 6. UPDATE STATE & RERUN (KUNCI PERBAIKAN TOMBOL DOWNLOAD)
                st.session_state.last_analysis_text = full_resp
                
                # PENTING: Rerun agar Sidebar membaca state terbaru dan menampilkan tombol
                st.rerun()

            except Exception as e:
                st.error(f"Terjadi kesalahan: {e}")
